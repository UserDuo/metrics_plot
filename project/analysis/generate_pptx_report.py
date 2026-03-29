# Installation:
# pip install python-pptx markdown-it-py

import os
from pptx import Presentation
from pptx.util import Inches
from markdown_it import MarkdownIt

def md_to_pptx(md_file, pptx_file):
    """
    Converts key sections of a Markdown file to a PPTX presentation.
    """
    md = MarkdownIt()
    
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()

    tokens = md.parse(md_content)
    
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Tokenizer Robustness Benchmark"
    subtitle.text = "An Information-Theoretic Analysis"

    # Sections to include
    sections_to_include = ["Executive Summary", "Key Findings", "Scientific Implications", "Recommendations"]
    
    current_section = ""
    content_buffer = ""

    for token in tokens:
        if token.type == 'heading_open' and token.tag == 'h2':
            # When a new section starts, add the previous one to a slide
            if current_section in sections_to_include and content_buffer:
                slide_layout = prs.slide_layouts[1] # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = current_section
                slide.shapes.placeholders[1].text = content_buffer.strip()
            
            content_buffer = ""
            current_section = ""

        elif token.type == 'inline' and token.level == 1:
             current_section = token.content
        
        elif current_section and token.content:
            content_buffer += token.content + "\n"

    # Add the last section
    if current_section in sections_to_include and content_buffer:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = current_section
        slide.shapes.placeholders[1].text = content_buffer.strip()

    prs.save(pptx_file)

if __name__ == "__main__":
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    md_path = os.path.join(project_root, "results", "tokenizer_benchmark_report.md")
    pptx_path = os.path.join(project_root, "results", "tokenizer_benchmark_report.pptx")

    print(f"Converting {md_path} to {pptx_path}...")
    md_to_pptx(md_path, pptx_path)
    print("Conversion complete.")
